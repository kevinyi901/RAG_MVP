"""Document loading with support for PDF, DOCX, and TXT."""

import os
import io
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, field
from pathlib import Path

# PDF processing
from pypdf import PdfReader

# DOCX processing
from docx import Document as DocxDocument

# PPTX processing
from pptx import Presentation

# Image processing
from PIL import Image

# Table extraction (optional)
try:
    import tabula
    TABULA_AVAILABLE = True
except ImportError:
    TABULA_AVAILABLE = False


@dataclass
class LoadedDocument:
    """Represents a loaded document with extracted text."""
    filename: str
    file_type: str
    file_size: int
    page_count: Optional[int]
    content: str
    pages: List[Dict[str, Any]]  # List of {page_number, content, section_title}
    tables: List[Dict[str, Any]] = field(default_factory=list)  # Extracted tables


class DocumentLoader:
    """Load and extract text from various document formats."""

    SUPPORTED_EXTENSIONS = {'.pdf', '.txt', '.md', '.docx', '.doc', '.pptx', '.png', '.jpg', '.jpeg'}

    def __init__(self):
        pass

    def load(self, file_path: str) -> LoadedDocument:
        """Load a document from file path."""
        path = Path(file_path)
        extension = path.suffix.lower()

        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")

        file_size = path.stat().st_size

        with open(file_path, 'rb') as f:
            content = f.read()

        return self.load_bytes(content, path.name, extension, file_size)

    def load_bytes(
        self,
        content: bytes,
        filename: str,
        file_type: str = None,
        file_size: int = None
    ) -> LoadedDocument:
        """Load a document from bytes."""
        if file_type is None:
            file_type = Path(filename).suffix.lower()

        if file_size is None:
            file_size = len(content)

        if file_type == '.pdf':
            return self._load_pdf(content, filename, file_size)
        elif file_type in ('.txt', '.md'):
            return self._load_text(content, filename, file_type, file_size)
        elif file_type in ('.docx', '.doc'):
            return self._load_docx(content, filename, file_size)
        elif file_type == '.pptx':
            return self._load_pptx(content, filename, file_size)
        elif file_type in ('.png', '.jpg', '.jpeg'):
            return self._load_image(content, filename, file_type, file_size)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _load_pdf(self, content: bytes, filename: str, file_size: int) -> LoadedDocument:
        """Load a PDF document."""
        reader = PdfReader(io.BytesIO(content))
        pages = []
        full_text = []

        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""

            # Try to extract section title from first line
            section_title = self._extract_section_title(page_text)
            
            # Simple extraction of images from PDF using pypdf if they exist
            images = []
            try:
                for img_obj in page.images:
                    images.append(img_obj.data)
            except Exception:
                pass

            pages.append({
                "page_number": page_num,
                "content": page_text,
                "section_title": section_title,
                "images": images
            })
            full_text.append(page_text)

        return LoadedDocument(
            filename=filename,
            file_type="pdf",
            file_size=file_size,
            page_count=len(reader.pages),
            content="\n\n".join(full_text),
            pages=pages
        )

    def _load_text(
        self,
        content: bytes,
        filename: str,
        file_type: str,
        file_size: int
    ) -> LoadedDocument:
        """Load a plain text or markdown file."""
        text = content.decode('utf-8', errors='replace')

        return LoadedDocument(
            filename=filename,
            file_type=file_type.lstrip('.'),
            file_size=file_size,
            page_count=None,
            content=text,
            pages=[{
                "page_number": None,
                "content": text,
                "section_title": None
            }]
        )

    def _load_docx(self, content: bytes, filename: str, file_size: int) -> LoadedDocument:
        """Load a DOCX document."""
        doc = DocxDocument(io.BytesIO(content))
        paragraphs = []
        current_section = None
        pages = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Check if this is a heading
            if para.style and para.style.name.startswith('Heading'):
                current_section = text

            paragraphs.append(text)

        full_text = "\n\n".join(paragraphs)

        # DOCX doesn't have real page numbers, treat as single page
        pages.append({
            "page_number": None,
            "content": full_text,
            "section_title": current_section
        })

        return LoadedDocument(
            filename=filename,
            file_type="docx",
            file_size=file_size,
            page_count=None,
            content=full_text,
            pages=pages
        )

    def _load_pptx(self, content: bytes, filename: str, file_size: int) -> LoadedDocument:
        """Load a PPTX document."""
        prs = Presentation(io.BytesIO(content))
        pages = []
        full_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_images = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
                elif shape.shape_type == 13: # 13 corresponds to msoPicture
                    # Extract image
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        slide_images.append(image_bytes)
                    except Exception:
                        pass
                        
            page_content = "\\n".join(t for t in slide_text if t)
            section_title = self._extract_section_title(page_content)
            
            pages.append({
                "page_number": i,
                "content": page_content,
                "section_title": section_title,
                "images": slide_images
            })
            full_text.append(page_content)

        return LoadedDocument(
            filename=filename,
            file_type="pptx",
            file_size=file_size,
            page_count=len(prs.slides),
            content="\\n\\n".join(full_text),
            pages=pages
        )

    def _load_image(self, content: bytes, filename: str, file_type: str, file_size: int) -> LoadedDocument:
        """Load an Image as a document."""
        return LoadedDocument(
            filename=filename,
            file_type=file_type.lstrip('.'),
            file_size=file_size,
            page_count=1,
            content=f"[Image File: {filename}]",
            pages=[{
                "page_number": 1,
                "content": f"[Image content for {filename}]",
                "section_title": None,
                "images": [content]
            }]
        )

    def _extract_section_title(self, text: str, max_length: int = 100) -> Optional[str]:
        """Extract a potential section title from text."""
        if not text:
            return None

        lines = text.strip().split('\n')
        if not lines:
            return None

        first_line = lines[0].strip()

        # Heuristics for section titles:
        # - Short lines (under max_length chars)
        # - All caps or title case
        # - Doesn't end with common sentence endings
        if len(first_line) > max_length:
            return None

        if first_line.endswith(('.', ',', ';', ':')):
            return None

        if first_line.isupper() or first_line.istitle():
            return first_line

        return None


class EnhancedDocumentLoader(DocumentLoader):
    """Enhanced loader with table extraction."""

    def __init__(
        self,
        extract_tables: bool = True,
        vllm_host: str = None,
        llm_model: str = None
    ):
        super().__init__()
        self.extract_tables = extract_tables
        self.vllm_host = vllm_host or os.getenv("VLLM_HOST", "http://localhost:8000")
        self.llm_model = llm_model or os.getenv("LLM_MODEL", "gpt-oss:20b")

    def _load_pdf(self, content: bytes, filename: str, file_size: int) -> LoadedDocument:
        """Load PDF with table extraction."""
        reader = PdfReader(io.BytesIO(content))
        pages = []
        full_text = []
        all_tables = []

        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""

            section_title = self._extract_section_title(page_text)

            # Simple extraction of images from PDF using pypdf if they exist
            images = []
            try:
                for img_obj in page.images:
                    images.append(img_obj.data)
            except Exception:
                pass

            pages.append({
                "page_number": page_num,
                "content": page_text,
                "section_title": section_title,
                "images": images
            })
            full_text.append(page_text)

        # Extract tables if enabled
        if self.extract_tables:
            all_tables = self._extract_tables_from_pdf(content)
            # Append table content to relevant pages
            for table in all_tables:
                table_text = self._table_to_text(table)
                if table.get("page"):
                    for p in pages:
                        if p["page_number"] == table["page"]:
                            p["content"] += f"\n\n[TABLE]\n{table_text}\n[/TABLE]"
                            break

        return LoadedDocument(
            filename=filename,
            file_type="pdf",
            file_size=file_size,
            page_count=len(reader.pages),
            content="\n\n".join(full_text),
            pages=pages,
            tables=all_tables
        )

    def _extract_tables_from_pdf(self, pdf_content: bytes) -> List[Dict[str, Any]]:
        """Extract tables from PDF using tabula if available."""
        if not TABULA_AVAILABLE:
            return []

        tables = []
        try:
            # Save content to temp file for tabula
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as f:
                f.write(pdf_content)
                temp_path = f.name

            # Extract tables
            dfs = tabula.read_pdf(temp_path, pages='all', multiple_tables=True)

            for i, df in enumerate(dfs):
                if df is not None and not df.empty:
                    tables.append({
                        "table_index": i,
                        "page": None,  # tabula doesn't reliably report page
                        "headers": list(df.columns),
                        "rows": df.values.tolist(),
                        "shape": df.shape
                    })

            # Clean up temp file
            os.unlink(temp_path)

        except Exception as e:
            print(f"Table extraction failed: {e}")

        return tables

    def _table_to_text(self, table: Dict[str, Any]) -> str:
        """Convert table dict to markdown text."""
        if not table.get("headers") or not table.get("rows"):
            return ""

        headers = table["headers"]
        rows = table["rows"]

        # Build markdown table
        lines = []
        lines.append("| " + " | ".join(str(h) for h in headers) + " |")
        lines.append("| " + " | ".join("---" for _ in headers) + " |")

        for row in rows:
            lines.append("| " + " | ".join(str(cell) for cell in row) + " |")

        return "\n".join(lines)
